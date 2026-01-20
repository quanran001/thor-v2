
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Configuration
MD_FILE = r'C:\01-工作区\03-sop-alchemist\docs\HACKATHON_PITCH.md'
OUTPUT_FILE = r'C:\01-工作区\03-sop-alchemist\docs\HACKATHON_PITCH_FINAL.pptx'
LOGO_PATH = r'C:\01-工作区\03-sop-alchemist\public\thor_avatar.png'

# Theme Colors (SOP Alchemist Style)
BG_COLOR = RGBColor(11, 15, 25)      # Slate-900 #0B0F19
ACCENT_COLOR = RGBColor(34, 211, 238) # Cyan-400 #22D3EE
TEXT_COLOR = RGBColor(248, 250, 252) # Slate-50 #F8FAFC
SEC_COLOR = RGBColor(148, 163, 184)   # Slate-400 #94A3B8

def parse_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by '---'
    raw_slides = content.split('\n---')
    slides_data = []
    
    for raw in raw_slides:
        if not raw.strip():
            continue
            
        slide = {
            'title': '',
            'screen_items': [],
            'notes': []
        }
        
        lines = raw.strip().split('\n')
        for line in lines:
            line = line.strip()
            # Match Title
            if '[Screen] 标题' in line or '[Screen] Title' in line:
                slide['title'] = line.split(':', 1)[1].strip()
            # Match Content
            elif '[Screen]' in line:
                text = line.split(':', 1)[1].strip() if ':' in line else line.replace('* [Screen]', '').strip()
                key = line.split(']')[1].split(':')[0].strip() if ':' in line else ""
                if key:
                    slide['screen_items'].append(f"{key}: {text}")
                else:
                    slide['screen_items'].append(text)
            # Match Bullets/Keywords (that are not labeled explicitly as Screen/Note but are bullets)
            elif line.startswith('*') or line.startswith('-'):
                clean = line[1:].strip()
                if '[Screen]' in clean: # Handle inline
                    continue # Handled above
                if '[Note]' in clean: # Handle inline note header
                    continue 
                # Determine if it belongs to Notes or Screen based on preceding lines?
                # Simple logic: If we hit a [Note] line, everything after is note.
                pass 
        
        # Re-parsing strategy: State Machine
        mode = 'screen' # or 'note'
        
        # Reset
        slide = { 'title': '', 'subtitle': '', 'bullets': [], 'notes': [] }
        
        for line in lines:
            line = line.strip()
            if not line or line.startswith('#'): continue
            
            if '[Note]' in line:
                mode = 'note'
                continue
            
            if '[Screen] 标题' in line:
                slide['title'] = line.split(':', 1)[1].strip()
            elif '[Screen] 副标题' in line:
                slide['subtitle'] = line.split(':', 1)[1].strip()
            elif '[Screen]' in line:
                # Other screen attributes like Slogan, Author
                clean_text = line.split(':', 1)[1].strip() if ':' in line else line
                # Add as bullet for now, or specific field
                slide['bullets'].append(clean_text)
            elif line.startswith('*') or line.startswith('-'):
                clean_bullet = line[1:].strip()
                # Skip empty bullets
                if not clean_bullet: continue
                
                if mode == 'note':
                    slide['notes'].append(clean_bullet)
                else:
                    # If it's a bullet in Screen section (like keywords)
                    # Use a heuristic: if we haven't seen [Note] yet, it's content
                    slide['bullets'].append(clean_bullet)

        slides_data.append(slide)
        
    return slides_data

def set_text_style(shape, size_pt, color, bold=False, align=None):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if align:
        tf.paragraphs[0].alignment = align
        
    for p in tf.paragraphs:
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(size_pt)
        p.font.color.rgb = color
        p.font.bold = bold

import random

# ... imports ...

def clean_markdown(text):
    # Remove **bold** markers but keep text
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Remove other markdown noise if any
    return text.strip()


BG_IMAGE_PATH = r'C:\01-工作区\03-sop-alchemist\public\space_bg.png'

def create_cool_ppt(slides_data):
    prs = Presentation()
    
    # 1. Setup Master Slide (Dark Mode)
    master = prs.slide_master
    bg = master.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR

    # Define Layouts
    # 0: Title Slide, 1: Title and Content, 6: Blank
    
    # --- Generate Slides ---
    for i, data in enumerate(slides_data):
        # Layout Decision
        if i == 0:
            layout = prs.slide_layouts[0] # Title Slide
        elif i == len(slides_data) - 1:
            layout = prs.slide_layouts[6] # Blank (Ending)
        else:
            layout = prs.slide_layouts[1] # Content
            
        slide = prs.slides.add_slide(layout)
        
        # --- Apply Background (Image) ---
        if os.path.exists(BG_IMAGE_PATH):
            # Add picture to cover 16:9 slide
            # W: 10 inches, H: 5.625 inches (default in python-pptx)
            # Or use prs.slide_width
            pic = slide.shapes.add_picture(BG_IMAGE_PATH, 0, 0, width=prs.slide_width, height=prs.slide_height)
            # Move to back (by moving to index 0 of spTree)
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element) # Index 0 or 2 depending on placeholder count
        else:
            # Fallback
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = BG_COLOR

        
        # --- Slide 1: Title Slide ---
        if i == 0:
            title = slide.shapes.title
            title.text = clean_markdown(data.get('title', 'SOP Alchemist'))
            
            subtitle = slide.placeholders[1]
            sub_text = clean_markdown(data.get('subtitle', ''))
            # Add other bullets to subtitle too
            for b in data['bullets']:
                sub_text += f"\n{clean_markdown(b)}"
            subtitle.text = sub_text
            
            # Styling
            set_text_style(title, 60, ACCENT_COLOR, True)
            set_text_style(subtitle, 24, TEXT_COLOR)
            
            # Add Logo
            if os.path.exists(LOGO_PATH):
                slide.shapes.add_picture(LOGO_PATH, Inches(4.5), Inches(3.0), width=Inches(1.5))

        # --- Middle Slides: Content ---
        elif 0 < i < len(slides_data) - 1:
            title = slide.shapes.title
            title.text = clean_markdown(data.get('title', 'Content'))
            set_text_style(title, 40, ACCENT_COLOR, True, PP_ALIGN.LEFT)
            
            # Body
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for bullet in data['bullets']:
                clean_bullet = clean_markdown(bullet)
                p = tf.add_paragraph()
                p.text = clean_bullet
                p.font.name = 'Microsoft YaHei'
                p.font.size = Pt(24)
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(14)
                
                # Simple Bold emulation: if line was fully bold in MD, make it bold?
                # For now, just clean.

            # Add decorative line
            line_shape = slide.shapes.add_shape(
                1, # MSO_SHAPE.RECTANGLE
                Inches(0.5), Inches(1.5), Inches(9), Inches(0.05)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = SEC_COLOR
            line_shape.line.fill.background() # No outline

        # --- Last Slide: Ending ---
        else:
            # Manually create Centered Text
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = clean_markdown(data.get('title', 'Thank You'))
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = ACCENT_COLOR
            p.alignment = PP_ALIGN.CENTER
            
            for b in data['bullets']:
                 p = tf.add_paragraph()
                 p.text = clean_markdown(b)
                 p.font.size = Pt(24)
                 p.font.color.rgb = TEXT_COLOR
                 p.alignment = PP_ALIGN.CENTER
        
        # --- Add Notes ---
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = "\n".join([clean_markdown(n) for n in data['notes']])
        
    prs.save(OUTPUT_FILE)
    print(f"PPT Generated: {OUTPUT_FILE}")


if __name__ == '__main__':
    try:
        data = parse_markdown(MD_FILE)
        create_cool_ppt(data)
    except Exception as e:
        print(f"Error: {e}")
